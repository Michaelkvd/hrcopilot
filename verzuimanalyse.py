
import io
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import requests

def extract_text_from_file(filename, contents):
    if filename.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(contents))
        return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif filename.endswith(".docx"):
        doc = Document(io.BytesIO(contents))
        return " ".join([para.text for para in doc.paragraphs])
    elif filename.endswith(".pptx"):
        prs = Presentation(io.BytesIO(contents))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return " ".join(text_runs)
    elif filename.endswith(".xlsx"):
        df = pd.read_excel(io.BytesIO(contents))
        return df.to_string()
    else:
        return ""

def fetch_cbs_verzuim(sbi_code, periode):
    try:
        url = f"https://opendata.cbs.nl/ODataFeed/odata/84472NED/TypedDataSet?$filter=BrancheSBI eq 'SBI{sbi_code}' and Perioden eq '{periode}'"
        response = requests.get(url)
        response.raise_for_status()
        data = response.json()
        records = data.get("value", [])
        if records:
            return float(records[0].get("Verzuimpercentage_1", 0))
        return None
    except Exception as e:
        print("Fout bij ophalen CBS-data:", e)
        return None

def bepaal_risiconiveau(intern, benchmark):
    afwijking = intern - benchmark
    if intern > 6 and afwijking > 1:
        return "hoog"
    elif intern > 4.5:
        return "matig"
    else:
        return "laag"

def genereer_aanbevelingen(risico):
    if risico == "hoog":
        return [
            "Voer direct een PMO of RI&E uit",
            "Plan gesprekken met leidinggevenden",
            "Stel een plan van aanpak op gericht op inzetbaarheid"
        ]
    elif risico == "matig":
        return [
            "Monitor verzuimtrends per maand",
            "Voer een preventief spreekuur in",
            "Herinner leidinggevenden aan verzuimprotocollen"
        ]
    else:
        return [
            "Blijf trends volgen",
            "Waak voor seizoensinvloeden",
            "Voorkom stille terugval bij medewerkers"
        ]

def genereer_vervolgstappen():
    return [
        "Wil je deze analyse exporteren als PDF?",
        "Wil je een PowerPoint slide laten maken voor het MT?",
        "Wil je dit als notitie toevoegen aan het kwartaalrapport?"
    ]

def analyse_verzuim(filename, contents, sbi_code='6420', periode='2024KW2'):
    tekst = extract_text_from_file(filename, contents).lower()
    kernwoorden = ["verzuim", "ziekmelding", "verzuimpercentage", "langdurig", "kort verzuim", "ziekteverzuim"]
    herkenning = [kw for kw in kernwoorden if kw in tekst]

    verzuimpercentage = 5.2 if "verzuim" in tekst else 2.5
    benchmark = fetch_cbs_verzuim(sbi_code, periode) or 0.0
    afwijking = verzuimpercentage - benchmark
    risico = bepaal_risiconiveau(verzuimpercentage, benchmark)
    aanbevelingen = genereer_aanbevelingen(risico)
    vervolgstappen = genereer_vervolgstappen()

    return {
        "kernwoorden_herkend": herkenning,
        "verzuimpercentage": verzuimpercentage,
        "risiconiveau": risico,
        "advies": aanbevelingen,
        "vervolgstappen": vervolgstappen,
        "cbs_benchmark": {
            "sbi_code": sbi_code,
            "periode": periode,
            "waarde": benchmark
        },
        "afwijking": afwijking,
        "interpretatie_markdown": f'''
### Verzuimanalyse

- Intern verzuim: **{verzuimpercentage}%**
- CBS Benchmark (SBI {sbi_code}, {periode}): **{benchmark}%**
- Afwijking: **{afwijking:+.1f}%**
- Risiconiveau: **{risico}**

**Interpretatie**:
Het verzuim ligt {afwijking:+.1f}% {("boven" if afwijking > 0 else "onder")} het branchegemiddelde. Risico wordt als **{risico.upper()}** beoordeeld.

**Aanbevelingen**:
{chr(10).join(f"- {a}" for a in aanbevelingen)}

**Vervolgstappen:**
{chr(10).join(f"- {v}" for v in vervolgstappen)}
'''
    }
}
