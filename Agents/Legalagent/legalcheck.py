
from fastapi import UploadFile
from typing import List, Optional, Tuple
import tempfile
import os
import extract_msg
from io import BytesIO
import pandas as pd
from pptx import Presentation

# Flexibele sleutelwoorden die de juridische module kunnen activeren.
TRIGGERS = {
    "juridisch",
    "legal",
    "ontslag",
    "contract",
    "vso",
    "vaststellingsovereenkomst",
    "arbeidsrecht",
    "beding",
}


def match_terms(text: str) -> bool:
    """Return ``True`` if juridische thema's vermoed worden in ``text``."""
    lower = text.lower()
    return any(t in lower for t in TRIGGERS)

KEYWORDS = [
    "ontslag", "verzuim", "vso", "vaststellingsovereenkomst", "brief",
    "officiële waarschuwing", "beëindiging", "dienstverband", "adviesaanvraag",
    "OVO", "overgang van onderneming", "beding", "leaver"
]
ALGEMENE_JURIDISCHE_BEGRIPPEN = [
    "concurrentiebeding", "relatiebeding", "studiekostenbeding", "arbeidsovereenkomst",
    "proeftijd", "eenzijdige wijziging", "transitievergoeding", "billijke vergoeding",
    "herplaatsing", "getuigschrift", "reorganisatie", "collectief ontslag",
    "disfunctioneren", "opzegtermijn", "op staande voet", "wederzijds goedvinden"
]

RISICO_NIVEAUS = {
    "concurrentiebeding": "hoog",
    "relatiebeding": "matig",
    "studiekostenbeding": "matig",
    "proeftijd": "laag",
    "transitievergoeding": "gemiddeld",
}

BRONNEN_INFO = {
    "wetten.nl": {
        "url": "https://wetten.overheid.nl",
        "omschrijving": "Volledige wetteksten ter onderbouwing van het juridisch kader.",
    },
    "rijksoverheid.nl": {
        "url": "https://www.rijksoverheid.nl",
        "omschrijving": "Officiële beleidsinformatie en toelichting van de overheid.",
    },
    "uwv.nl": {
        "url": "https://www.uwv.nl",
        "omschrijving": "Regelingen rondom werk, verzuim en uitkeringen.",
    },
    "ontslag.nl": {
        "url": "https://www.ontslag.nl",
        "omschrijving": "Achtergrondinformatie over beëindiging van arbeidsovereenkomsten.",
    },
    "maxius.nl": {
        "url": "https://maxius.nl",
        "omschrijving": "Geannoteerde wetgeving en jurisprudentie.",
    },
}

def extract_text_from_input(
    file: Optional[UploadFile] = None,
    input_text: Optional[str] = None
) -> str:
    if input_text and len(input_text) > 20:
        return input_text
    if file:
        fname = file.filename.lower()
        try:
            if fname.endswith((".pdf", ".docx", ".doc", ".txt", ".eml")):
                return file.file.read().decode("utf-8", errors="ignore")

            if fname.endswith(".msg"):
                with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as tmp:
                    tmp.write(file.file.read())
                    tmp.flush()
                    msg = extract_msg.Message(tmp.name)
                    text = f"{msg.subject}\n{msg.body}" if msg.body else msg.subject
                os.unlink(tmp.name)
                return text or ""

            if fname.endswith((".xlsx", ".xls", ".csv")):
                buf = BytesIO(file.file.read())
                buf.seek(0)
                try:
                    df = pd.read_excel(buf)
                except Exception:
                    buf.seek(0)
                    df = pd.read_csv(buf)
                return " ".join(df.astype(str).stack().tolist())

            if fname.endswith((".ppt", ".pptx")):
                buf = BytesIO(file.file.read())
                prs = Presentation(buf)
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slides_text.append(shape.text)
                return "\n".join(slides_text)

            return file.file.read().decode("utf-8", errors="ignore")
        except Exception:
            return ""
    return ""

def flexibele_begrippenherkenning(text: str) -> Tuple[list, list]:
    gevonden_kernwoorden = [word for word in KEYWORDS if word.lower() in text.lower()]
    gevonden_juridische = [begrip for begrip in ALGEMENE_JURIDISCHE_BEGRIPPEN if begrip.lower() in text.lower()]
    return gevonden_kernwoorden, gevonden_juridische

def casus_complexiteit_score(keywords: list, juridische_begrippen: list) -> str:
    totaal = len(keywords) + len(juridische_begrippen)
    if totaal == 0:
        return "eenvoudig"
    elif totaal <= 2:
        return "middelmatig"
    else:
        return "complex"

def risico_inschatting(
    text: str, complexiteit: str, keywords: list, juridische_begrippen: list
) -> str:
    """Bepaal het risiconiveau op een flexibele manier."""

    score = len(text) / 300
    score += len(juridische_begrippen)
    # Weeg complexiteit zwaarder mee dan eventuele kernwoorden
    score += {"eenvoudig": 0, "middelmatig": 1.5, "complex": 3}[complexiteit]
    score += len(keywords) * 0.5
    score += len(juridische_begrippen)
    score += {"eenvoudig": 0, "middelmatig": 1, "complex": 2}[complexiteit]
    if score >= 6:
        return "hoog"
    if score >= 3:
        return "matig"
    return "laag"

def bronnen_check(text: str, juridische_begrippen: list, complexiteit: str) -> dict:
    """Selecteer relevante bronnen en geef omschrijving en link terug."""

    import re
    from collections import Counter

    woorden = re.findall(r"[a-zA-Z]{4,}", text.lower())
    onderwerp = Counter(woorden).most_common(1)[0][0] if woorden else "algemeen"
    artikel = f"art-{abs(hash(onderwerp)) % 100}"

    bronnen: dict[str, list[tuple[str, str]]] = {}

    info = BRONNEN_INFO["wetten.nl"]
    bronnen["wetten.nl"] = [(f"{info['url']}/{artikel}", info["omschrijving"])]

    lc = text.lower()

    if "verzuim" in lc or "ziekte" in lc or complexiteit != "eenvoudig":
        info = BRONNEN_INFO["uwv.nl"]
        bronnen["uwv.nl"] = [(f"{info['url']}/{artikel}", info["omschrijving"])]

    if "ontslag" in lc or "be\u00ebindiging" in lc:
        info = BRONNEN_INFO["ontslag.nl"]
        bronnen["ontslag.nl"] = [(f"{info['url']}/{artikel}", info["omschrijving"])]

    if complexiteit != "eenvoudig" or len(juridische_begrippen) > 1:
        info = BRONNEN_INFO["rijksoverheid.nl"]
        bronnen["rijksoverheid.nl"] = [(f"{info['url']}/{artikel}", info["omschrijving"])]

    if juridische_begrippen:
        info = BRONNEN_INFO["maxius.nl"]
        bronnen["maxius.nl"] = [(f"{info['url']}/{artikel}", info["omschrijving"])]

    return bronnen
  

def genereer_vragen(kernwoorden: List[str], juridische_begrippen: List[str]) -> List[str]:
    vragen: List[str] = []
    if "ontslag" in kernwoorden:
        vragen.append("Welke stappen zijn al ondernomen richting ontslag?")
    if "verzuim" in kernwoorden:
        vragen.append("Is het verzuimprotocol volledig gevolgd?")
    if "concurrentiebeding" in juridische_begrippen:
        vragen.append("Bestaat er een geldig concurrentiebeding in het contract?")
    if not vragen:
        vragen.append("Kunt u extra informatie delen over de situatie?")
    return vragen

def generate_legal_advice(
    kernwoorden: list,
    juridische_begrippen: list,
    bronnen: dict,
    complexiteit: str,
    input_text: str,
    intern_beleid: Optional[str] = None
) -> Tuple[str, str, List[str], str]:
    # Flexibele GPT-stijl adviezen
    if not kernwoorden and not juridische_begrippen:
        advies = (
            "Op basis van de aangeleverde informatie zijn er geen duidelijke juridische risico’s gedetecteerd. "
            "Toch is het verstandig om alert te blijven op signalen die alsnog tot juridische vragen kunnen leiden."
        )
        actieplan = (
            "Voor nu zijn er geen directe acties vereist. "
            "Mocht de situatie veranderen of aanvullende informatie beschikbaar komen, heroverweeg dan deze analyse."
        )
        vragen = genereer_vragen(kernwoorden, juridische_begrippen)
        risico = "laag"
        return advies, actieplan, vragen, risico

    # Dynamisch advies op basis van context
    aandachtspunten = kernwoorden + juridische_begrippen
    advies = (
        f"Deze casus bevat juridische aandachtspunten rondom: {', '.join(aandachtspunten)}. "
        f"{'Het betreft een complexe zaak waarbij zorgvuldig handelen essentieel is.' if complexiteit == 'complex' else ''}"
        " Overweeg welke belangen er spelen, raadpleeg waar nodig een jurist en neem voldoende tijd om alle relevante stappen te doorlopen. "
    )
    if intern_beleid:
        advies += f"Neem aanvullend het interne beleid in acht: {intern_beleid} "

    # Dynamisch actieplan
    actieplan = "Een mogelijk vervolgtraject bestaat uit:\n"
    stappen = []
    stappen.append("• Leg alle communicatie en acties volledig vast in het dossier.")
    if complexiteit == "complex":
        stappen.append("• Schakel direct juridische ondersteuning of een HR-expert in.")
    else:
        stappen.append("• Beoordeel of externe ondersteuning noodzakelijk is of dat interne afstemming volstaat.")
    if bronnen:
        stappen.append("• Raadpleeg relevante wet- en regelgeving, bijvoorbeeld:")
        for bron, hits in bronnen.items():
            for wet, uitleg in hits:
                stappen.append(f"   - {bron}: {wet} – {uitleg}")
    else:
        stappen.append("• Controleer of aanvullende bronnen geraadpleegd moeten worden.")
    stappen.append("• Evalueer de situatie regelmatig en pas het plan waar nodig aan.")

    actieplan += "\n".join(stappen)
    risico = risico_inschatting(input_text, complexiteit, kernwoorden, juridische_begrippen)
    vragen = genereer_vragen(kernwoorden, juridische_begrippen)
    return advies, actieplan, vragen, risico

def legalcheck(
    file: Optional[UploadFile] = None,
    input_text: Optional[str] = None,
    intern_beleid: Optional[str] = None
) -> dict:
    text = extract_text_from_input(file=file, input_text=input_text)
    if not text or len(text) < 20:
        return {
            "status": "onvoldoende input",
            "advies": (
                "De aangeleverde input is te beperkt om een juridische beoordeling te kunnen geven. "
                "Lever meer context, een document of aanvullende informatie aan voor een grondige analyse."
            ),
            "actieplan": (
                "Voeg een relevante brief, beleidsstuk of extra context toe zodat een gerichte analyse mogelijk is."
            ),
            "legal_markdown": "## Juridische Analyse\n**Status:** Onvoldoende input voor analyse.\n"
        }

    kernwoorden, juridische_begrippen = flexibele_begrippenherkenning(text)
    complexiteit = casus_complexiteit_score(kernwoorden, juridische_begrippen)
    bronnen = bronnen_check(text, juridische_begrippen, complexiteit)
    advies, actieplan, vragen, risico = generate_legal_advice(
        kernwoorden, juridische_begrippen, bronnen, complexiteit, text, intern_beleid
    )

    markdown = f"""## Juridische Analyse
**Herkenbare kernwoorden:** {', '.join(kernwoorden) or 'Geen'}
**Herkenbare juridische begrippen:** {', '.join(juridische_begrippen) or 'Geen'}
**Complexiteit casus:** {complexiteit}
**Advies:**
{advies}
**Actieplan:**
{actieplan}
**Relevante bronnen/artikelen:**
"""
    for bron, info_list in bronnen.items():
        for url, oms in info_list:
            markdown += f"- **{bron}**: {oms} ({url})\n"
    if not bronnen:
        markdown += "Geen specifieke bronnen gevonden.\n"
    markdown += f"\n**Risico-inschatting:** {risico}\n"
    if intern_beleid:
        markdown += f"\n**Interne beleidsinformatie:**\n{intern_beleid}\n"
    if vragen:
        markdown += "\n**Vervolgvragen:**\n"
        for q in vragen:
            markdown += f"- {q}\n"

    return {
        "status": "ok",
        "herkenning_kernwoorden": kernwoorden,
        "herkenning_juridische_begrippen": juridische_begrippen,
        "complexiteit": complexiteit,
        "advies": advies,
        "actieplan": actieplan,
        "bronnen": bronnen,
        "verdiepende_vragen": vragen,
        "risico": risico,
        "legal_markdown": markdown
    }
