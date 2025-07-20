from fastapi import UploadFile
from typing import Optional
from utils import text_matches
from io import BytesIO
import pandas as pd
import extract_msg
from pptx import Presentation
import re
import tempfile
import os

# Triggers that indicate compliance related questions
TRIGGERS = {
    "compliance",
    "avg",
    "gdpr",
    "privacy",
    "persoonsgegevens",
    "beveiliging",
    "security",
    "policy",
    "datalek",
    "data breach",
    "data protection",
}


def match_terms(text: str) -> bool:
    """Return True if compliance gerelateerde termen voorkomen in ``text``."""
    return text_matches(text, TRIGGERS)


def extract_text(file: Optional[UploadFile] = None, text: Optional[str] = None) -> str:
    if text:
        return text
    if not file:
        return ""
    fname = file.filename.lower()
    try:
        if fname.endswith((".pdf", ".docx", ".doc", ".txt", ".eml")):
            return file.file.read().decode("utf-8", errors="ignore")
        if fname.endswith(".msg"):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as tmp:
                tmp.write(file.file.read())
                tmp.flush()
                msg = extract_msg.Message(tmp.name)
                text = f"{msg.subject}\n{msg.body}" if msg.body else msg.subject
            os.unlink(tmp.name)
            return text or ""
        if fname.endswith((".xlsx", ".xls", ".csv")):
            buf = BytesIO(file.file.read())
            buf.seek(0)
            try:
                df = pd.read_excel(buf)
            except Exception:
                buf.seek(0)
                df = pd.read_csv(buf)
            return " ".join(df.astype(str).stack().tolist())
        if fname.endswith((".ppt", ".pptx")):
            buf = BytesIO(file.file.read())
            prs = Presentation(buf)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return "\n".join(slides_text)
        return file.file.read().decode("utf-8", errors="ignore")
    except Exception:
        return ""


def detect_pii(text: str) -> list[str]:
    emails = re.findall(r"[\w.-]+@[\w.-]+", text)
    numbers = re.findall(r"\b\d{8,}\b", text)
    return list(set(emails + numbers))


def compliance_check(file: Optional[UploadFile] = None, text: Optional[str] = None) -> dict:
    content = extract_text(file=file, text=text)
    if not content:
        return {"status": "geen input", "issues": []}
    lc = content.lower()
    hits = [kw for kw in TRIGGERS if kw in lc]
    pii = detect_pii(content)
    advies = "Controleer document op naleving van privacy- en beveiligingsrichtlijnen."
    status = "ok" if hits or pii else "geen bijzonderheden"
    return {
        "status": status,
        "gevonden_termen": hits,
        "gevonden_pii": pii,
        "advies": advies,
    }
